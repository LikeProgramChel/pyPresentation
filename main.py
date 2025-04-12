import sys
import numpy as np
from pptx import Presentation
from pptx.util import Pt, Inches
import random
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout,
                             QHBoxLayout, QPushButton, QTextEdit, QLineEdit,
                             QLabel, QListWidget, QMessageBox, QFileDialog)
from PyQt6.QtCore import Qt, QMimeData, QPoint, QRect
from PyQt6.QtGui import QPixmap, QDragEnterEvent, QDropEvent, QMouseEvent, QPainter, QFont
import os


# Класс нейронной сети
class PresentationAgent:
    def __init__(self, input_size=5, hidden_size=10, output_size=3):
        self.W1 = np.random.randn(input_size, hidden_size) * 0.01
        self.W2 = np.random.randn(hidden_size, output_size) * 0.01
        self.learning_rate = 0.01

    def sigmoid(self, x):
        return 1 / (1 + np.exp(-x))

    def forward(self, state):
        self.z1 = np.dot(state, self.W1)
        self.a1 = self.sigmoid(self.z1)
        self.z2 = np.dot(self.a1, self.W2)
        self.a2 = self.sigmoid(self.z2)
        return self.a2

    def train(self, state, reward, action):
        output = self.forward(state)
        error = reward - output[action]
        delta2 = error * self.a2 * (1 - self.a2)
        delta1 = np.dot(delta2, self.W2.T) * self.a1 * (1 - self.a1)
        self.W2 += self.learning_rate * np.outer(self.a1, delta2)
        self.W1 += self.learning_rate * np.outer(state, delta1)


# Класс для создания презентации
class PresentationMaker:
    def __init__(self):
        self.agent = PresentationAgent()
        self.templates = [
            {"title_size": Pt(40), "text_size": Pt(24), "layout": 0},
            {"title_size": Pt(36), "text_size": Pt(20), "layout": 1},
            {"title_size": Pt(44), "text_size": Pt(28), "layout": 2}
        ]

    def create_slide(self, title, content, slide_num):
        state = np.array([
            slide_num / 10.0,
            len(title) / 50.0,
            len(content) / 200.0,
            slide_num % 3,
            random.random()
        ])
        action_probs = self.agent.forward(state)
        action = np.argmax(action_probs)
        return action

    def make_presentation(self, slides_data, filename="presentation.pptx"):
        prs = Presentation()
        previous_action = 0
        logo_path = "logo.png"

        for i, (title, content, image_data) in enumerate(slides_data):
            action = self.create_slide(title, content, i)
            template = self.templates[action]
            slide_layout = prs.slide_layouts[template["layout"]]
            slide = prs.slides.add_slide(slide_layout)

            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            title_placeholder.text_frame.paragraphs[0].font.size = template["title_size"]

            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = content
                content_placeholder.text_frame.paragraphs[0].font.size = template["text_size"]

            if image_data:
                image_path, left, top = image_data
                try:
                    slide.shapes.add_picture(image_path, left, top, width=Inches(4))
                except FileNotFoundError:
                    print(f"Не удалось загрузить изображение: {image_path}")

            if os.path.exists(logo_path):
                try:
                    logo_width = Inches(1)
                    left = Inches(10) - logo_width - Inches(0.2)
                    top = Inches(7.5) - logo_width - Inches(0.2)
                    slide.shapes.add_picture(logo_path, left, top, width=logo_width)
                except FileNotFoundError:
                    print("Не удалось загрузить logo.png")
            else:
                print("Файл logo.png не найден в текущей директории")

            reward = random.uniform(0, 1)
            state = np.array([i / 10.0, len(title) / 50.0, len(content) / 200.0, previous_action, random.random()])
            self.agent.train(state, reward, action)
            previous_action = action

        prs.save(filename)
        return f"Презентация сохранена как {filename}"


# Кастомный QLabel для предпросмотра слайда
class SlidePreviewLabel(QLabel):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setMouseTracking(True)
        self.dragging = False
        self.offset = QPoint()
        self.image_position = QPoint(0, 0)
        self.current_slide = None
        self.logo_pixmap = QPixmap("logo.png") if os.path.exists("logo.png") else QPixmap()

    def set_slide(self, slide_data):
        self.current_slide = slide_data
        if slide_data and slide_data[2]:
            _, _, (image_path, left, top) = slide_data
            pixmap = QPixmap(image_path)
            if not pixmap.isNull():
                scaled_pixmap = pixmap.scaled(120, 80, Qt.AspectRatioMode.KeepAspectRatio)
                self.setPixmap(scaled_pixmap)
                x = int((left / Inches(10)) * 300)
                y = int((top / Inches(7.5)) * 200)
                self.image_position = QPoint(x, y)
            else:
                self.setPixmap(QPixmap())
        else:
            self.setPixmap(QPixmap())
        self.update()

    def mousePressEvent(self, event: QMouseEvent):
        if event.button() == Qt.MouseButton.LeftButton and not self.pixmap().isNull():
            self.dragging = True
            self.offset = event.pos() - self.image_position

    def mouseMoveEvent(self, event: QMouseEvent):
        if self.dragging:
            new_pos = event.pos() - self.offset
            new_pos.setX(max(0, min(new_pos.x(), self.width() - self.pixmap().width())))
            new_pos.setY(max(0, min(new_pos.y(), self.height() - self.pixmap().height())))
            self.image_position = new_pos
            self.update()

    def mouseReleaseEvent(self, event: QMouseEvent):
        if event.button() == Qt.MouseButton.LeftButton:
            self.dragging = False
            if self.current_slide and self.current_slide[2]:
                title, content, (image_path, _, _) = self.current_slide
                left = Inches((self.image_position.x() / 300) * 10)
                top = Inches((self.image_position.y() / 200) * 7.5)
                self.current_slide = (title, content, (image_path, left, top))

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.fillRect(self.rect(), Qt.GlobalColor.white)

        if self.current_slide:
            title, content, image_data = self.current_slide

            painter.setFont(QFont("Arial", 12, QFont.Weight.Bold))
            painter.drawText(QRect(10, 10, 280, 30), Qt.AlignmentFlag.AlignLeft | Qt.TextFlag.TextWordWrap, title)

            painter.setFont(QFont("Arial", 8))
            painter.drawText(QRect(10, 50, 280, 100), Qt.AlignmentFlag.AlignLeft | Qt.TextFlag.TextWordWrap, content)

            if not self.pixmap().isNull():
                painter.drawPixmap(self.image_position, self.pixmap())

            if not self.logo_pixmap.isNull():
                logo_width = 30
                scaled_logo = self.logo_pixmap.scaled(logo_width, logo_width, Qt.AspectRatioMode.KeepAspectRatio)
                painter.drawPixmap(300 - logo_width - 6, 200 - logo_width - 6, scaled_logo)

        painter.setPen(Qt.GlobalColor.black)
        painter.drawRect(0, 0, 299, 199)


# Графический интерфейс
class PresentationApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("pyPresentation")
        self.setGeometry(100, 100, 800, 500)

        self.maker = PresentationMaker()
        self.slides_data = []

        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        layout = QHBoxLayout(main_widget)

        # Левая часть: ввод данных
        left_layout = QVBoxLayout()

        self.title_input = QLineEdit()
        self.title_input.setPlaceholderText("Введите заголовок слайда")
        self.content_input = QTextEdit()
        self.content_input.setPlaceholderText("Введите текст слайда")
        self.image_label = QLabel("Изображение не выбрано")
        image_button = QPushButton("Выбрать изображение")
        image_button.clicked.connect(self.select_image)

        button_layout = QHBoxLayout()
        add_button = QPushButton("Добавить слайд")
        add_button.clicked.connect(self.add_slide)
        create_button = QPushButton("Создать презентацию")
        create_button.clicked.connect(self.create_presentation)
        open_button = QPushButton("Открыть проект")
        open_button.clicked.connect(self.open_project)
        button_layout.addWidget(add_button)
        button_layout.addWidget(create_button)
        button_layout.addWidget(open_button)

        self.slide_list = QListWidget()
        self.slide_list.itemDoubleClicked.connect(self.remove_slide)
        self.slide_list.itemClicked.connect(self.show_preview)

        left_layout.addWidget(self.title_input)
        left_layout.addWidget(self.content_input)
        left_layout.addWidget(self.image_label)
        left_layout.addWidget(image_button)
        left_layout.addLayout(button_layout)
        left_layout.addWidget(QLabel("Дважды кликните на слайд, чтобы удалить"))
        left_layout.addWidget(self.slide_list)

        # Правая часть: предпросмотр
        right_layout = QVBoxLayout()
        self.preview_label = SlidePreviewLabel(self)
        self.preview_label.setFixedSize(300, 200)
        self.preview_label.setStyleSheet("border: 1px solid #ccc;")
        right_layout.addWidget(self.preview_label)
        right_layout.addStretch()

        layout.addLayout(left_layout)
        layout.addLayout(right_layout)

        self.setAcceptDrops(True)
        self.current_image_path = None
        self.current_left = Inches(1)
        self.current_top = Inches(2)

        self.setStyleSheet("""
            QLineEdit, QTextEdit {
                border: 1px solid #ccc;
                border-radius: 5px;
                padding: 5px;
            }
            QPushButton {
                background-color: #4CAF50;
                color: white;
                padding: 5px;
                border-radius: 5px;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
            QLabel {
                padding: 2px;
            }
        """)

    def select_image(self):
        image_path, _ = QFileDialog.getOpenFileName(self, "Выбрать изображение", "",
                                                    "Images (*.png *.jpg *.jpeg *.bmp)")
        if image_path and os.path.basename(image_path) != "logo.png":
            self.current_image_path = image_path
            self.image_label.setText(f"Выбрано: {os.path.basename(image_path)}")
            self.current_left = Inches(1)
            self.current_top = Inches(2)
            self.update_preview(image_path)
        elif os.path.basename(image_path) == "logo.png":
            QMessageBox.warning(self, "Ошибка", "Нельзя использовать logo.png!")
        else:
            self.current_image_path = None
            self.image_label.setText("Изображение не выбрано")
            self.preview_label.set_slide(None)

    def add_slide(self):
        title = self.title_input.text().strip()
        content = self.content_input.toPlainText().strip()
        if title and content:
            if self.current_image_path:
                left = Inches((self.preview_label.image_position.x() / 300) * 10)
                top = Inches((self.preview_label.image_position.y() / 200) * 7.5)
                image_data = (self.current_image_path, left, top)
            else:
                image_data = None
            self.slides_data.append((title, content, image_data))
            self.slide_list.addItem(f"Слайд {len(self.slides_data)}: {title}")
            self.title_input.clear()
            self.content_input.clear()
            self.current_image_path = None
            self.image_label.setText("Изображение не выбрано")
            self.preview_label.set_slide(None)
        else:
            QMessageBox.warning(self, "Ошибка", "Введите заголовок и текст!")

    def remove_slide(self, item):
        index = self.slide_list.row(item)
        self.slides_data.pop(index)
        self.slide_list.takeItem(index)
        self.preview_label.set_slide(None)

    def create_presentation(self):
        if not self.slides_data:
            QMessageBox.warning(self, "Ошибка", "Добавьте хотя бы один слайд!")
            return
        result = self.maker.make_presentation(self.slides_data)
        QMessageBox.information(self, "Успех", result)

    def show_preview(self, item):
        index = self.slide_list.row(item)
        self.preview_label.set_slide(self.slides_data[index])

    def update_preview(self, image_path):
        slide_data = (self.title_input.text(), self.content_input.toPlainText(),
                      (image_path, self.current_left, self.current_top) if image_path else None)
        self.preview_label.set_slide(slide_data)

    def open_project(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Открыть проект", "", "PowerPoint Files (*.pptx)")
        if file_path:
            try:
                prs = Presentation(file_path)
                self.slides_data.clear()
                self.slide_list.clear()
                self.preview_label.set_slide(None)

                for i, slide in enumerate(prs.slides):
                    title = slide.shapes.title.text if slide.shapes.title else f"Слайд {i + 1}"
                    content = ""
                    image_data = None

                    # Извлечение текста из второго placeholder, если он есть
                    if len(slide.placeholders) > 1:
                        content = slide.placeholders[1].text

                    # Поиск изображения (кроме логотипа)
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            image_path = shape.image.filename if hasattr(shape.image, 'filename') else None
                            if image_path and os.path.basename(image_path) != "logo.png" and os.path.exists(image_path):
                                left = shape.left
                                top = shape.top
                                image_data = (image_path, left, top)
                            elif image_path and os.path.basename(image_path) != "logo.png":
                                # Если файл изображения не найден, сохраняем путь, но предупреждаем
                                print(f"Изображение {image_path} не найдено в текущей директории")
                                image_data = (image_path, shape.left, shape.top)

                    self.slides_data.append((title, content, image_data))
                    self.slide_list.addItem(f"Слайд {len(self.slides_data)}: {title}")

                QMessageBox.information(self, "Успех", f"Проект {os.path.basename(file_path)} загружен!")
            except Exception as e:
                QMessageBox.warning(self, "Ошибка", f"Не удалось открыть проект: {str(e)}")

    def dragEnterEvent(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls():
            url = event.mimeData().urls()[0]
            file_path = url.toLocalFile()
            if file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp')) and os.path.basename(
                    file_path) != "logo.png":
                event.acceptProposedAction()

    def dropEvent(self, event: QDropEvent):
        file_path = event.mimeData().urls()[0].toLocalFile()
        if file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp')) and os.path.basename(file_path) != "logo.png":
            self.current_image_path = file_path
            self.image_label.setText(f"Выбрано: {os.path.basename(file_path)}")
            self.current_left = Inches(1)
            self.current_top = Inches(2)
            self.update_preview(file_path)
        else:
            QMessageBox.warning(self, "Ошибка", "Нельзя использовать logo.png или неподдерживаемый формат!")


# Запуск приложения
if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = PresentationApp()
    window.show()
    sys.exit(app.exec())